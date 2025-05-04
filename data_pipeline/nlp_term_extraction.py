import spacy
import pdfplumber
from docx import Document
from pptx import Presentation
from pathlib import Path
import os
import json

# Load spaCy model
nlp = spacy.load("en_core_web_sm")

def extract_text(file_path):
    path = Path(file_path)
    ext = path.suffix.lower()

    if ext == ".txt":
        return path.read_text(encoding="utf-8", errors="ignore")

    elif ext == ".docx":
        doc = Document(path)
        return "\n".join(para.text for para in doc.paragraphs)

    elif ext == ".pptx":
        prs = Presentation(path)
        text_runs = []
        for slide in prs.slides:
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    text_runs.append(shape.text)
        return "\n".join(text_runs)

    elif ext == ".pdf":
        text = ""
        with pdfplumber.open(path) as pdf:
            for page in pdf.pages:
                page_text = page.extract_text()
                if page_text:
                    text += page_text + "\n"
        return text

    else:
        raise ValueError(f"Unsupported file type: {ext}")

def extract_entities_from_file(file_path):
    try:
        text = extract_text(file_path)
        doc = nlp(text)
        entities = sorted(set(
            ent.text.strip() for ent in doc.ents
            if ent.label_ in {"PERSON", "ORG", "GPE", "DATE"}
        ))
        return entities
    except Exception as e:
        print(f"Error processing {file_path}: {e}")
        return []

def preview_entities(base_dir="data/Processed_Committees", limit=50):
    count = 0
    for root, _, files in os.walk(base_dir):
        for file in files:
            if file.endswith(".json"):
                json_path = os.path.join(root, file)
                base_name = Path(file).stem
                for ext in [".txt", ".docx", ".pptx", ".pdf"]:
                    doc_path = os.path.join(root, base_name + ext)
                    if os.path.exists(doc_path):
                        entities = extract_entities_from_file(doc_path)
                        print(f"{base_name}{ext}: {entities}")
                        count += 1
                        break
                if count >= limit:
                    return

def run_entity_preview(base_dir="data/Processed_Committees", limit=25):
    preview_entities(base_dir=base_dir, limit=limit)

if __name__ == "__main__":
    run_entity_preview()
